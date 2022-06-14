# import required module
import os
import fitz
import docx2txt
import textract
import docx
import re
from pptx import Presentation


def add_bookmark(paragraph, bookmark_text, bookmark_name):
    run = paragraph.add_run()
    # for reference the following also works: tag =  document.element.xpath('//w:r')[-1]
    tag = run._r
    start = docx.oxml.shared.OxmlElement('w:bookmarkStart')
    start.set(docx.oxml.ns.qn('w:id'), '0')
    start.set(docx.oxml.ns.qn('w:name'), bookmark_name)
    tag.append(start)

    text = docx.oxml.OxmlElement('w:r')
    text.text = bookmark_text
    tag.append(text)

    end = docx.oxml.shared.OxmlElement('w:bookmarkEnd')
    end.set(docx.oxml.ns.qn('w:id'), '0')
    end.set(docx.oxml.ns.qn('w:name'), bookmark_name)
    tag.append(end)


def find_all(regex, texttosearch):
    match_list = []
    while True:
        match = re.search(regex, texttosearch, re.IGNORECASE)
        if match:
            # match_list.append(match.group(0))
            match_list.append(str(match.span(0)))
            texttosearch = texttosearch[match.end():]
        else:
            return match_list

# assign directory


def searchPdf(file_in_dir, filename):
    # counter for hits
    hitcount = 0
    pdf_file = fitz.open(file_in_dir)
    # insert document details as header for the results
    # separator line
    print('<hr>', file=f)
    print('<h3 title><a href="' + file_in_dir +
          '" target="_blank">' + filename + '</a></h3>', file=f)
    print('<h3 title>' + 'PDF File ' + '</h3>', file=f)
    # print('<h4><b>Author:</b> ' + str(df['Author'][ind]) + '</h4>', file=f)
    print('<h4>' + str(pdf_file.pageCount) + ' pages</h4>', file=f)
    for pageNumber, page in enumerate(pdf_file.pages(), start=1):
        pdf_text = page.get_text()
        ResSearch = find_all(string, pdf_text)
        if len(ResSearch) != 0:
            hitcount += 1
            # separator line
            print('<hr>', file=f)
            print('<h5>Page ' + str(pageNumber) + ' of ' + str(pdf_file.pageCount) + '  <a href="' + file_in_dir +
                  '#page=' + str(pageNumber) + '" target="_blank">Link</a></h5>', file=f)        # page number
            # use for offset for 1+ matches
            counter = 0
            for hit in ResSearch:
                comma = int(hit.find(','))
                startchar = int(hit[1:comma]) + counter
                endchar = int(hit[comma + 2:len(hit) - 1]) + counter
                print('<hr>', file=f)
                # separator line
                # the returned text is 'buffertext' characters before and after the search term, and the search term is highlighted in italics
                # the <i> tag is formatted in the CSS section of the HTML file
                buffertext = pdf_text[max(startchar-textbuffer, 0):startchar] + '<i>' + \
                    pdf_text[startchar:endchar] + '</i>' + \
                    pdf_text[endchar:endchar+textbuffer]
                buffertext = ''.join(buffertext.splitlines())
                print('<p>' + buffertext + '</p>', file=f)
                counter = endchar

    if hitcount == 0:
        # separator line
        print('<hr>', file=f)
        print('<p>No hits</p>', file=f)
    # print('\n', pageNumber)


def searchDoc(file_in_dir, filename):
    # counter for hits
    hitcount = 0
    docx_text = docx.Document(file_in_dir)
    num_para = len(docx_text.paragraphs)
    # insert document details as header for the results
    # separator line
    print('<hr>', file=f)
    print('<h3 title><a href="' + file_in_dir +
          '" target="_blank">' + filename + '</a></h3>', file=f)
    print('<h3 title>' + 'DOCX File ' + '</h3>', file=f)
    # print('<h4><b>Author:</b> ' + str(df['Author'][ind]) + '</h4>', file=f)
    print('<h4>' + str(num_para) + ' Paragraphs</h4>', file=f)

    curr_para = 1
    for paranum, para in enumerate(docx_text.paragraphs):
        doc_text = para.text
        ResSearch = find_all(string, doc_text)
        if len(ResSearch) != 0:
            # add a bookmakr to every paragraph
            # add_bookmark(paragraph=para, bookmark_text=f"temp{paranum}", bookmark_name=f"temp{paranum+1}")
            hitcount += 1
            # separator line
            print('<hr>', file=f)
            print('<h5>Paragraph ' + str(curr_para) + ' of ' + str(num_para) + '  <a href="' + file_in_dir +
                   '#Bookmark=' + f'temp{paranum}' + '" target="_blank">Link</a></h5>', file=f)        # page number
            # use for offset for 1+ matches
            counter = 0
            for hit in ResSearch:
                comma = int(hit.find(','))
                startchar = int(hit[1:comma]) + counter
                endchar = int(hit[comma + 2:len(hit) - 1]) + counter
                print('<hr>', file=f)
                # separator line
                # the returned text is 'buffertext' characters before and after the search term, and the search term is highlighted in italics
                # the <i> tag is formatted in the CSS section of the HTML file
                buffertext = doc_text[max(startchar-textbuffer, 0):startchar] + '<i>' + \
                    doc_text[startchar:endchar] + '</i>' + \
                    doc_text[endchar:endchar+textbuffer]
                buffertext = ''.join(buffertext.splitlines())
                print('<p>' + buffertext + '</p>', file=f)
                counter = endchar
        curr_para += 1
    # docx_text.save(file_in_dir)
    if hitcount == 0:
        # separator line
        print('<hr>', file=f)
        print('<p>No hits</p>', file=f)

def searchPptx(file_in_dir, filename):
    # counter for hits
    hitcount = 0
    pptx_text = Presentation(file_in_dir)
    num_of_slides = len(pptx_text.slides)
    # insert document details as header for the results
    # separator line
    print('<hr>', file=f)
    print('<h3 title><a href="' + file_in_dir +
          '" target="_blank">' + filename + '</a></h3>', file=f)
    print('<h3 title>' + 'PPTX File ' + '</h3>', file=f)
    # print('<h4><b>Author:</b> ' + str(df['Author'][ind]) + '</h4>', file=f)
    print('<h4>' + str(num_of_slides) + ' slides</h4>', file=f)
    for idx, slide in enumerate(pptx_text.slides):
        pptx_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                pptx_text += shape.text
        ResSearch = find_all(string, pptx_text)
        # print('slide number %d' % idx+1)
        if len(ResSearch) != 0:
            # add a bookmakr to every paragraph
            # add_bookmark(paragraph=para, bookmark_text=f"temp{paranum}", bookmark_name=f"temp{paranum+1}")
            hitcount += 1
            # separator line
            print('<hr>', file=f)
            print('<h5>Slide ' + str(idx+1) + ' of ' + str(num_of_slides) + '  <a href="' + file_in_dir +
                   '#Slide=' + str(idx+1) + '" target="_blank">Link</a></h5>', file=f)        # page number
            # use for offset for 1+ matches
            counter = 0
            for hit in ResSearch:
                comma = int(hit.find(','))
                startchar = int(hit[1:comma]) + counter
                endchar = int(hit[comma + 2:len(hit) - 1]) + counter
                print('<hr>', file=f)
                # separator line
                # the returned text is 'buffertext' characters before and after the search term, and the search term is highlighted in italics
                # the <i> tag is formatted in the CSS section of the HTML file
                buffertext = pptx_text[max(startchar-textbuffer, 0):startchar] + '<i>' + \
                    pptx_text[startchar:endchar] + '</i>' + \
                    pptx_text[endchar:endchar+textbuffer]
                buffertext = ''.join(buffertext.splitlines())
                print('<p>' + buffertext + '</p>', file=f)
                counter = endchar
    # docx_text.save(file_in_dir)
    if hitcount == 0:
        # separator line
        print('<hr>', file=f)
        print('<p>No hits</p>', file=f)


# regular expression search string
string = 'data|information|delete'
# main directory containing files
directory = 'C:\TMP\Python_PDF_Search'
# characters to return from searched document before/after a search hit
textbuffer = 300
# output file to save search results
searchresults = 'C:/TMP/Python_PDF_Search/Search_Results(' + string.replace(
    '|', ' OR ') + ').html'
# open file to save results
f = open(searchresults, 'w', encoding="utf-8")
print('<html>', file=f)
print('<head>', file=f)

# css for HTML file to make the result outlook look pretty
print('<style>', file=f)
print('p    {font-family: segoe ui;font-size: 80%;}', file=f)
print('h3   {font-family: segoe ui;font-size: 120%;}', file=f)
print('h4   {font-family: segoe ui;font-size: 100%;}', file=f)
print('i    {font-family: segoe ui;font-size: 100%;color: crimson;}', file=f)
print('</style>', file=f)

print('</head>', file=f)
print('<body>', file=f)

print('<h1>PDF_Search_HTML_Output_EXAMPLE.py</h1>', file=f)
print('<h2>Search results for <i>' +
      string.replace('|', '</i> OR <i>') + '</i></h2>', file=f)

# -------------------------------------------------------------------------------------------------------------------------------------------
# iterate over files in
# that directory
for filename in os.listdir(directory):
    file_in_dir = os.path.join(directory, filename)
    # checking if it is a file
    if os.path.isfile(file_in_dir):
        name, extension = os.path.splitext(file_in_dir)
        if extension == '.pdf':
            searchPdf(file_in_dir, filename)
        elif extension == '.docx':
            searchDoc(file_in_dir, filename)
        elif extension == '.pptx':
            searchPptx(file_in_dir, filename)
        else:
            print('file not supported yet')
        # print(extension)
        # print(f)

#-------------------------------------------------------------------------------------------------------------------------------------------
print('<hr>', file = f)                                                             # separator line
print('</body>', file = f)
print('</html>', file = f)
f.close()                                                                           # close the results file
#-------------------------------------------------------------------------------------------------------------------------------------------
